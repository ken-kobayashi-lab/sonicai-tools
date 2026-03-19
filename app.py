
import io
import os
from copy import copy
from datetime import date

import pandas as pd
import plotly.graph_objects as go
import streamlit as st
from openpyxl import load_workbook
from openpyxl.styles import Alignment
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# ==========================================
# 0. パス定義
# ==========================================
current_dir = os.path.dirname(__file__)
logo_path = os.path.join(current_dir, "logo.png")
template_path = os.path.join(current_dir, "template.pptx")
quotation_template_path = os.path.join(current_dir, "quotation_template.xlsx")

# ==========================================
# 1. session_state 初期化
# ==========================================
def init_session_state():
    defaults = {
        "my_company": "株式会社SonicAI",
        "my_zip": "108-0075",
        "my_address": "東京都港区港南2-16-1 7F Spaces品川",
        "my_name": "小林賢正",
        "my_mail": "ken-kobayashi@sonicai.jp",
        "my_tel": "080-8044-3236",
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


DEFAULT_PRODUCT_MASTER = [
    {"商品名": "AI検査装置A", "単位": "式", "単価": 500000, "備考": ""},
    {"商品名": "AI検査装置B", "単位": "式", "単価": 800000, "備考": ""},
    {"商品名": "カメラ", "単位": "台", "単価": 120000, "備考": ""},
    {"商品名": "レンズ", "単位": "個", "単価": 50000, "備考": ""},
    {"商品名": "照明", "単位": "台", "単価": 80000, "備考": ""},
    {"商品名": "制御PC", "単位": "台", "単価": 200000, "備考": ""},
    {"商品名": "導入設定費", "単位": "式", "単価": 150000, "備考": ""},
]

DETAIL_START_ROW = 19
DETAIL_END_ROW = 27
MAX_DETAIL_ROWS = DETAIL_END_ROW - DETAIL_START_ROW + 1


def init_quotation_state():
    if "product_master_df" not in st.session_state:
        st.session_state["product_master_df"] = pd.DataFrame(DEFAULT_PRODUCT_MASTER)

    if "quotation_rows" not in st.session_state:
        st.session_state["quotation_rows"] = [
            {
                "マスタ選択": "自由入力",
                "品名": "",
                "仕様・備考": "",
                "数量": 1,
                "単位": "",
                "単価": 0,
                "行値引き": 0,
            }
            for _ in range(5)
        ]


# Streamlit UI作成前に必ず初期化
init_session_state()
init_quotation_state()

# ==========================================
# 2. 共通関数
# ==========================================
def fill_placeholder(slide, ph_type_idx, content, is_image=False, img_file=None):
    placeholders = sorted(
        [sh for sh in slide.shapes if sh.is_placeholder],
        key=lambda x: (x.top, x.left)
    )

    text_phs = [
        sh for sh in placeholders
        if sh.placeholder_format.type != PP_PLACEHOLDER.PICTURE
    ]
    img_phs = [
        sh for sh in placeholders
        if sh.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    ]

    try:
        if is_image:
            if img_file and len(img_phs) > ph_type_idx:
                ph = img_phs[ph_type_idx]
                ph.insert_picture(img_file)
                return True
            return False
        else:
            if len(text_phs) > ph_type_idx:
                ph = text_phs[ph_type_idx]
                ph.text = str(content)
                return True
            return False
    except Exception:
        return False


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slide_id = xml_slides[old_index]
    xml_slides.remove(slide_id)
    xml_slides.insert(new_index, slide_id)


def safe_int(value):
    try:
        if value is None or value == "":
            return 0
        return int(round(float(value)))
    except Exception:
        return 0


def format_yen(value):
    try:
        return f"¥{int(round(float(value))):,}"
    except Exception:
        return "¥0"


def write_wrapped_text(ws, cell_ref, value):
    ws[cell_ref] = value
    if ws[cell_ref].alignment:
        ws[cell_ref].alignment = copy(ws[cell_ref].alignment)
        ws[cell_ref].alignment = Alignment(
            horizontal=ws[cell_ref].alignment.horizontal,
            vertical="top",
            text_rotation=ws[cell_ref].alignment.textRotation,
            wrap_text=True,
            shrink_to_fit=ws[cell_ref].alignment.shrinkToFit,
            indent=ws[cell_ref].alignment.indent,
        )
    else:
        ws[cell_ref].alignment = Alignment(vertical="top", wrap_text=True)


def clear_detail_row(ws, row_idx):
    for col in ["C", "D", "E", "F", "G", "H", "I"]:
        ws[f"{col}{row_idx}"] = None


# ==========================================
# 3. Report Generator ロジック
# ==========================================
def validate_report_data(data):
    required_keys = [
        "company_name",
        "project_name",
        "my_company",
        "my_zip",
        "my_address",
        "my_name",
        "my_mail",
        "my_tel",
        "cond",
        "setup_img",
        "res",
        "uploaded_images",
        "image_comments",
        "summ",
    ]
    missing = [k for k in required_keys if k not in data]
    if missing:
        raise ValueError(f"必要なデータが不足しています: {missing}")


def generate_pptx(template_path_arg, data):
    validate_report_data(data)

    if not os.path.exists(template_path_arg):
        raise FileNotFoundError(f"テンプレートファイルが見つかりません: {template_path_arg}")

    prs = Presentation(template_path_arg)

    fill_placeholder(prs.slides[0], 0, data["company_name"])
    fill_placeholder(prs.slides[0], 1, data["project_name"])

    creator_info = (
        f"{data['my_company']}\n"
        f"〒{data['my_zip']}\n"
        f"{data['my_address']}\n"
        f"{data['my_name']}\n"
        f"mail：{data['my_mail']}\n"
        f"TEL：{data['my_tel']}"
    )
    fill_placeholder(prs.slides[0], 2, creator_info)

    fill_placeholder(prs.slides[1], 1, data["cond"])
    if data["setup_img"] is not None:
        fill_placeholder(prs.slides[1], 0, "", is_image=True, img_file=data["setup_img"])

    fill_placeholder(prs.slides[2], 1, data["res"])

    images = data["uploaded_images"]
    comments = data["image_comments"]

    if len(prs.slides) < 4:
        raise ValueError("template.pptx の4枚目（画像貼り付け用スライド）が不足しています。")

    img_layout = prs.slides[3].slide_layout

    for i in range(0, len(images), 2):
        slide = prs.slides[3] if i == 0 else prs.slides.add_slide(img_layout)

        if i > 0:
            move_slide(prs, len(prs.slides) - 1, 3 + (i // 2))

        fill_placeholder(slide, 0, "", is_image=True, img_file=images[i])
        if len(comments) > i:
            fill_placeholder(slide, 1, comments[i])

        if i + 1 < len(images):
            fill_placeholder(slide, 1, "", is_image=True, img_file=images[i + 1])
            if len(comments) > i + 1:
                fill_placeholder(slide, 2, comments[i + 1])
        else:
            fill_placeholder(slide, 2, "")

    fill_placeholder(prs.slides[-1], 1, data["summ"])

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io


# ==========================================
# 4. Quotation Generator ロジック
# ==========================================
def apply_product_master_to_row(row_data, selected_name, product_master_df):
    if selected_name == "自由入力":
        row_data["マスタ選択"] = "自由入力"
        return row_data

    matched = product_master_df[product_master_df["商品名"] == selected_name]
    if matched.empty:
        row_data["マスタ選択"] = "自由入力"
        return row_data

    record = matched.iloc[0]
    row_data["マスタ選択"] = selected_name
    row_data["品名"] = record.get("商品名", "")
    row_data["単位"] = record.get("単位", "")
    row_data["単価"] = safe_int(record.get("単価", 0))
    if not row_data.get("仕様・備考"):
        row_data["仕様・備考"] = record.get("備考", "")
    return row_data


def calculate_row_amount(row):
    qty = safe_int(row.get("数量", 0))
    unit_price = safe_int(row.get("単価", 0))
    line_discount = safe_int(row.get("行値引き", 0))
    return max(qty * unit_price - line_discount, 0)


def generate_quotation_excel(data):
    if not os.path.exists(quotation_template_path):
        raise FileNotFoundError(
            f"quotation_template.xlsx が見つかりません: {quotation_template_path}"
        )

    wb = load_workbook(quotation_template_path)
    if "見積書" not in wb.sheetnames:
        raise ValueError("quotation_template.xlsx に『見積書』シートがありません。")

    ws = wb["見積書"]

    ws["C6"] = data["customer_name"]
    ws["H6"] = data["quote_no"]
    ws["C7"] = data["subject"]
    ws["H7"] = data["quote_date_str"]

    ws["C9"] = data["payment_terms"]
    ws["C10"] = data["delivery_terms"]
    ws["C11"] = data["valid_until"]

    ws["G10"] = data["my_company"]
    ws["G11"] = f"〒{data['my_zip']}"
    ws["G12"] = data["my_address"]
    ws["G13"] = data["my_name"]
    ws["G14"] = f"mail: {data['my_mail']}"
    ws["G15"] = f"TEL: {data['my_tel']}"

    for r in range(DETAIL_START_ROW, DETAIL_END_ROW + 1):
        clear_detail_row(ws, r)

    detail_rows = data["detail_rows"][:MAX_DETAIL_ROWS]
    for idx, item in enumerate(detail_rows):
        row_no = DETAIL_START_ROW + idx
        row_amount = calculate_row_amount(item)

        ws[f"C{row_no}"] = item.get("品名", "")
        write_wrapped_text(ws, f"D{row_no}", item.get("仕様・備考", ""))
        ws[f"F{row_no}"] = safe_int(item.get("数量", 0))
        ws[f"G{row_no}"] = item.get("単位", "")
        ws[f"H{row_no}"] = safe_int(item.get("単価", 0))
        ws[f"I{row_no}"] = row_amount

    subtotal = safe_int(data["subtotal_before_global_discount"])
    global_discount = safe_int(data["global_discount"])
    subtotal_after_discount = max(subtotal - global_discount, 0)
    tax_amount = safe_int(data["tax_amount"])
    total_including_tax = safe_int(data["total_including_tax"])

    ws["I28"] = subtotal_after_discount
    ws["I29"] = tax_amount
    ws["I30"] = total_including_tax

    # 税抜合計表示
    ws["C15"] = subtotal_after_discount

    remarks_lines = [
        f"支払条件：{data['payment_terms']}" if data["payment_terms"] else "",
        f"納期：{data['delivery_terms']}" if data["delivery_terms"] else "",
        f"有効期限：{data['valid_until']}" if data["valid_until"] else "",
        f"全体値引き：{format_yen(global_discount)}" if global_discount else "",
        data["remarks"] if data["remarks"] else "",
    ]
    remarks_text = "\n".join([x for x in remarks_lines if x])
    write_wrapped_text(ws, "C33", remarks_text)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output


# ==========================================
# 5. Streamlit 基本設定
# ==========================================
st.set_page_config(page_title="SonicAI AI Tools", layout="wide")

st.markdown("""
<style>
    html, body, [class*="css"] {
        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
    }
    .logo-container {
        text-align: center;
        padding-bottom: 20px;
    }
    div.stButton > button {
        background-color: #ff6600;
        color: white;
        height: 3em;
        width: 100%;
        font-weight: bold;
        border-radius: 10px;
        border: none;
        transition: 0.3s;
    }
    div.stButton > button:hover {
        background-color: #e65c00;
        box-shadow: 0 4px 15px rgba(255,102,0,0.3);
    }
    div.stDownloadButton > button {
        background-color: #28a745;
        color: white;
        width: 100%;
        height: 3em;
        font-weight: bold;
        border-radius: 10px;
        border: none;
    }
    .camera-info-card {
        background-color: #1e2630;
        color: #ffffff;
        padding: 20px;
        border-radius: 12px;
        border-left: 6px solid #007bff;
        margin: 10px 0;
        line-height: 1.6;
    }
    .camera-info-label {
        color: #8899ac;
        font-size: 0.85em;
        font-weight: bold;
        text-transform: uppercase;
    }
    .camera-info-value {
        color: #ffffff;
        font-size: 1.1em;
        margin-bottom: 8px;
    }
    .section-card {
        background: #f8f9fb;
        border: 1px solid #e9edf2;
        border-radius: 12px;
        padding: 16px;
        margin-bottom: 12px;
    }
</style>
""", unsafe_allow_html=True)

# ==========================================
# 6. サイドバー
# ==========================================
with st.sidebar:
    st.markdown('<div class="logo-container">', unsafe_allow_html=True)
    if os.path.exists(logo_path):
        st.image(logo_path, use_container_width=True)
    else:
        st.title("SonicAI")
    st.markdown('</div>', unsafe_allow_html=True)

    st.title("SonicAI Platform")
    current_tab = st.radio(
        "Navigation",
        ["Tools", "Settings", "Manual"],
        label_visibility="collapsed"
    )

    tool_choice = None
    if current_tab == "Tools":
        tool_choice = st.selectbox(
            "Select Tool",
            ["📄 Report Generator", "🔍 Optical Calc", "💰 Quotation Generator"]
        )

    elif current_tab == "Settings":
        st.subheader("👤 Creator Info")
        st.text_input("Company", key="my_company")
        st.text_input("Zip Code", key="my_zip")
        st.text_area("Address", key="my_address", height=80)
        st.text_input("Name", key="my_name")
        st.text_input("Email", key="my_mail")
        st.text_input("TEL", key="my_tel")
        st.caption("Settings are used for report / quotation output.")

    st.write("---")
    st.caption("SonicAI Inc. v1.8")

# ==========================================
# 7. メインコンテンツ
# ==========================================
if current_tab == "Manual":
    st.markdown("<h1 style='color: #6c757d;'>📖 User Manual</h1>", unsafe_allow_html=True)

    col_m1, col_m2 = st.columns(2)

    with col_m1:
        st.subheader("🔍 視野計算ツールの使い方")
        st.markdown("""
1. **カメラ・レンズ選択**: 使用するカメラ型式とレンズの焦点距離を選びます。  
2. **モード選択**: WD（距離）から視野を出すか、必要な視野からWDを出すかを選びます。  
3. **グラフの活用**: オレンジ色の線が選択中のレンズ特性です。赤い点線が現在の計算値を示します。  
        """)

    with col_m2:
        st.subheader("📄 レポート作成ツールの使い方")
        st.markdown("""
1. **基本情報入力**: 案件名や検証条件（照明設定など）を入力します。  
2. **画像アップロード**: 検証に使用した画像をドラッグ＆ドロップします。  
3. **生成**: 「パワーポイントを生成」ボタンを押し、完成したファイルをダウンロードします。  
        """)

        st.subheader("💰 見積もり作成ツールの使い方")
        st.markdown("""
1. **基本情報入力**: 客先名、件名、見積番号、支払条件などを入力します。  
2. **商品マスタ or 自由入力**: 明細ごとに商品を選ぶか、自由入力で記載します。  
3. **値引き設定**: 行値引きと全体値引きを入力します。  
4. **生成**: 「見積書Excelを生成」ボタンでテンプレートに反映されたファイルをダウンロードします。  
        """)

    st.divider()
    st.warning("⚠️ 本ツールはオフライン動作を前提としています。ライブラリ更新時以外はネット接続不要です。")

elif current_tab == "Settings":
    st.markdown("<h1 style='color: #6c757d;'>👤 Settings Preview</h1>", unsafe_allow_html=True)
    st.write("サイドバーで入力した情報は、以下の形式で出力ファイルに反映されます。")
    st.code(
        f"{st.session_state.get('my_company', '株式会社SonicAI')}\n"
        f"〒{st.session_state.get('my_zip', '108-0075')}\n"
        f"{st.session_state.get('my_address', '東京都港区港南2-16-1 7F Spaces品川')}\n"
        f"{st.session_state.get('my_name', '小林賢正')}\n"
        f"mail：{st.session_state.get('my_mail', 'ken-kobayashi@sonicai.jp')}\n"
        f"TEL：{st.session_state.get('my_tel', '080-8044-3236')}"
    )

elif tool_choice == "📄 Report Generator":
    st.markdown("<h1 style='color: #ff6600;'>📄 Report Generator</h1>", unsafe_allow_html=True)

    c_name = st.text_input("客先会社名", value="〇〇株式会社 御中")
    p_name = st.text_input("案件名", value="AI外観検査 導入可否検証")

    st.write("---")
    col1, col2 = st.columns(2)

    with col1:
        cond = st.text_area("検証条件", value="カメラ：\nレンズ：\n照明：", height=150)
        setup_img = st.file_uploader("構成図(任意)", type=["png", "jpg", "jpeg"])

    with col2:
        res = st.text_area("検証結果", value="良好な結果を確認しました。", height=150)
        summ = st.text_area("ご提案まとめ", value="本構成での導入を推奨します。", height=150)

    st.write("---")

    uploaded_images = st.file_uploader(
        "画像をアップロード",
        accept_multiple_files=True,
        type=["png", "jpg", "jpeg"]
    )

    image_comments = []
    if uploaded_images:
        cols = st.columns(2)
        for idx, img in enumerate(uploaded_images):
            with cols[idx % 2]:
                st.image(img, use_container_width=True)
                comment = st.text_input(
                    f"Comment {idx + 1}",
                    value=f"画像({idx + 1}): {img.name}",
                    key=f"c_{idx}"
                )
                image_comments.append(comment)

    if st.button("🚀 パワーポイントを生成", use_container_width=True):
        try:
            if not uploaded_images:
                st.error("画像を1枚以上選択してください。")
            else:
                data = {
                    "company_name": c_name,
                    "project_name": p_name,
                    "my_company": st.session_state.get("my_company", "株式会社SonicAI"),
                    "my_zip": st.session_state.get("my_zip", "108-0075"),
                    "my_address": st.session_state.get("my_address", "東京都港区港南2-16-1 7F Spaces品川"),
                    "my_name": st.session_state.get("my_name", "小林賢正"),
                    "my_mail": st.session_state.get("my_mail", "ken-kobayashi@sonicai.jp"),
                    "my_tel": st.session_state.get("my_tel", "080-8044-3236"),
                    "cond": cond,
                    "setup_img": setup_img,
                    "res": res,
                    "uploaded_images": uploaded_images,
                    "image_comments": image_comments,
                    "summ": summ,
                }

                final_ppt = generate_pptx(template_path, data)

                st.success("パワーポイントを生成しました。")
                st.balloons()
                st.download_button(
                    label="📥 ダウンロード",
                    data=final_ppt,
                    file_name=f"{p_name}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        except Exception as e:
            st.error(f"エラーが発生しました: {e}")

elif tool_choice == "💰 Quotation Generator":
    st.markdown("<h1 style='color: #28a745;'>💰 Quotation Generator</h1>", unsafe_allow_html=True)

    if not os.path.exists(quotation_template_path):
        st.error("quotation_template.xlsx が見つかりません。app.py と同じフォルダに配置してください。")
        st.stop()

    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader("基本情報")
    q_col1, q_col2, q_col3 = st.columns(3)

    with q_col1:
        customer_name = st.text_input("客先会社名", value="〇〇株式会社 御中")
        subject = st.text_input("件名", value="AI外観検査装置 お見積り")
        quote_no = st.text_input("見積番号", value="SONIC-2026-001")

    with q_col2:
        quote_date = st.date_input("見積日", value=date.today())
        valid_until = st.text_input("有効期限", value="発行日より30日")
        payment_terms = st.text_input("支払条件", value="月末締め翌月末払い")

    with q_col3:
        delivery_terms = st.text_input("納期", value="ご注文後6〜8週間")
        tax_rate_percent = st.number_input("消費税率(%)", min_value=0.0, max_value=100.0, value=10.0, step=0.1)
        global_discount = st.number_input("全体値引き(円)", min_value=0, value=0, step=1000)

    remarks = st.text_area("備考", value="御見積条件等ございましたらご連絡ください。", height=100)
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader("商品マスタ")
    st.caption("ここで編集した商品マスタは、このセッション中の見積入力に反映されます。")
    product_master_df = st.data_editor(
        st.session_state["product_master_df"],
        num_rows="dynamic",
        use_container_width=True,
        key="product_master_editor",
        column_config={
            "商品名": st.column_config.TextColumn("商品名", required=True),
            "単位": st.column_config.TextColumn("単位"),
            "単価": st.column_config.NumberColumn("単価", min_value=0, step=1000, format="%d"),
            "備考": st.column_config.TextColumn("備考"),
        }
    )
    st.session_state["product_master_df"] = product_master_df.copy()
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader(f"明細入力（最大 {MAX_DETAIL_ROWS} 行）")

    detail_row_count = st.number_input(
        "入力行数",
        min_value=1,
        max_value=MAX_DETAIL_ROWS,
        value=min(len(st.session_state["quotation_rows"]), MAX_DETAIL_ROWS),
        step=1
    )

    while len(st.session_state["quotation_rows"]) < detail_row_count:
        st.session_state["quotation_rows"].append(
            {
                "マスタ選択": "自由入力",
                "品名": "",
                "仕様・備考": "",
                "数量": 1,
                "単位": "",
                "単価": 0,
                "行値引き": 0,
            }
        )

    st.session_state["quotation_rows"] = st.session_state["quotation_rows"][:detail_row_count]

    product_names = ["自由入力"]
    if not product_master_df.empty and "商品名" in product_master_df.columns:
        product_names += [str(x) for x in product_master_df["商品名"].fillna("").tolist() if str(x).strip()]

    detail_rows = []
    for i in range(detail_row_count):
        base = st.session_state["quotation_rows"][i]

        with st.container():
            st.markdown(f"**明細 {i + 1}**")
            c1, c2, c3, c4, c5, c6, c7 = st.columns([1.4, 2.0, 1.2, 1.0, 1.2, 1.2, 1.3])

            selected_master = c1.selectbox(
                "マスタ",
                options=product_names,
                index=product_names.index(base.get("マスタ選択", "自由入力")) if base.get("マスタ選択", "自由入力") in product_names else 0,
                key=f"q_master_{i}"
            )

            row_data = {
                "マスタ選択": selected_master,
                "品名": base.get("品名", ""),
                "仕様・備考": base.get("仕様・備考", ""),
                "数量": safe_int(base.get("数量", 1)),
                "単位": base.get("単位", ""),
                "単価": safe_int(base.get("単価", 0)),
                "行値引き": safe_int(base.get("行値引き", 0)),
            }

            row_data = apply_product_master_to_row(row_data, selected_master, product_master_df)

            row_data["品名"] = c2.text_input("品名", value=row_data["品名"], key=f"q_name_{i}")
            row_data["数量"] = c3.number_input("数量", min_value=0, value=safe_int(row_data["数量"]), step=1, key=f"q_qty_{i}")
            row_data["単位"] = c4.text_input("単位", value=row_data["単位"], key=f"q_unit_{i}")
            row_data["単価"] = c5.number_input("単価", min_value=0, value=safe_int(row_data["単価"]), step=1000, key=f"q_price_{i}")
            row_data["行値引き"] = c6.number_input("行値引き", min_value=0, value=safe_int(row_data["行値引き"]), step=1000, key=f"q_discount_{i}")
            row_amount = calculate_row_amount(row_data)
            c7.metric("金額", format_yen(row_amount))

            row_data["仕様・備考"] = st.text_input(
                "仕様・備考",
                value=row_data["仕様・備考"],
                key=f"q_spec_{i}"
            )

            detail_rows.append(row_data)
            st.session_state["quotation_rows"][i] = row_data
            st.write("---")
    st.markdown('</div>', unsafe_allow_html=True)

    subtotal_before_global_discount = sum(calculate_row_amount(r) for r in detail_rows)
    subtotal_after_global_discount = max(subtotal_before_global_discount - safe_int(global_discount), 0)
    tax_amount = int(round(subtotal_after_global_discount * (tax_rate_percent / 100)))
    total_including_tax = subtotal_after_global_discount + tax_amount

    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader("集計")
    s1, s2, s3, s4 = st.columns(4)
    s1.metric("小計（行値引き反映後）", format_yen(subtotal_before_global_discount))
    s2.metric("全体値引き後税抜", format_yen(subtotal_after_global_discount))
    s3.metric("消費税", format_yen(tax_amount))
    s4.metric("合計（税込）", format_yen(total_including_tax))
    st.markdown('</div>', unsafe_allow_html=True)

    preview_df = pd.DataFrame([
        {
            "品名": r.get("品名", ""),
            "仕様・備考": r.get("仕様・備考", ""),
            "数量": safe_int(r.get("数量", 0)),
            "単位": r.get("単位", ""),
            "単価": safe_int(r.get("単価", 0)),
            "行値引き": safe_int(r.get("行値引き", 0)),
            "金額": calculate_row_amount(r),
        }
        for r in detail_rows
        if str(r.get("品名", "")).strip() or safe_int(r.get("数量", 0)) > 0 or safe_int(r.get("単価", 0)) > 0
    ])

    st.subheader("プレビュー")
    if preview_df.empty:
        st.info("明細を入力するとここにプレビューが表示されます。")
    else:
        st.dataframe(preview_df, use_container_width=True, hide_index=True)

    if st.button("📥 見積書Excelを生成", use_container_width=True):
        try:
            quote_data = {
                "customer_name": customer_name,
                "subject": subject,
                "quote_no": quote_no,
                "quote_date_str": quote_date.strftime("%Y/%m/%d"),
                "payment_terms": payment_terms,
                "delivery_terms": delivery_terms,
                "valid_until": valid_until,
                "remarks": remarks,
                "detail_rows": detail_rows,
                "subtotal_before_global_discount": subtotal_before_global_discount,
                "global_discount": global_discount,
                "tax_amount": tax_amount,
                "total_including_tax": total_including_tax,
                "my_company": st.session_state.get("my_company", "株式会社SonicAI"),
                "my_zip": st.session_state.get("my_zip", "108-0075"),
                "my_address": st.session_state.get("my_address", "東京都港区港南2-16-1 7F Spaces品川"),
                "my_name": st.session_state.get("my_name", "小林賢正"),
                "my_mail": st.session_state.get("my_mail", "ken-kobayashi@sonicai.jp"),
                "my_tel": st.session_state.get("my_tel", "080-8044-3236"),
            }

            generated_excel = generate_quotation_excel(quote_data)
            file_base = customer_name.replace(" ", "_").replace("/", "_")
            st.success("見積書を生成しました。")
            st.download_button(
                label="📥 ダウンロード",
                data=generated_excel,
                file_name=f"{file_base}_見積書.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        except Exception as e:
            st.error(f"見積書生成時にエラーが発生しました: {e}")

else:
    st.markdown("<h1 style='color: #007bff;'>🔍 Optical Calc Tool</h1>", unsafe_allow_html=True)

    camera_data = {
        "40万画素 (CS-41-C/B)": {"model": "CS-41-C/B", "sensor": "1/2.9", "w": 720, "h": 540, "px": 6.9},
        "160万画素 (CS-160-C/B)": {"model": "CS-160-C/B", "sensor": "1/2.9", "w": 1440, "h": 1080, "px": 3.45},
        "500万画素 (CS-500-C/B)": {"model": "CS-500-C/B", "sensor": "1/2", "w": 2600, "h": 2160, "px": 2.5}
    }
    focal_lengths = [8, 12, 16, 25, 35, 50, 75]

    c1, c2 = st.columns([1, 1])

    with c1:
        cam_label = st.selectbox("Select Camera", list(camera_data.keys()))
        cam = camera_data[cam_label]
        st.markdown(f"""
        <div class="camera-info-card">
            <div class="camera-info-label">Model</div>
            <div class="camera-info-value">{cam['model']}</div>
            <div class="camera-info-label">Sensor Size</div>
            <div class="camera-info-value">{cam['sensor']}</div>
            <div class="camera-info-label">Pixel Size</div>
            <div class="camera-info-value">{cam['px']}μm × {cam['px']}μm</div>
            <div class="camera-info-label">Resolution</div>
            <div class="camera-info-value">{cam['w']} × {cam['h']} px</div>
        </div>
        """, unsafe_allow_html=True)

    with c2:
        focal_len = st.selectbox("Focal Length (mm)", focal_lengths)
        f_value = st.number_input("F-Value", value=4.0, step=0.1)
        calc_mode = st.radio("Mode", ["WD -> FOV", "FOV -> WD"])

    sensor_w = cam["w"] * cam["px"] / 1000

    if calc_mode == "WD -> FOV":
        wd_val = st.number_input("Working Distance (mm)", value=300.0, step=10.0)
        fov_w = wd_val * sensor_w / focal_len
    else:
        fov_w = st.number_input("Target FOV (Width mm)", value=100.0, step=10.0)
        wd_val = focal_len * fov_w / sensor_w

    fig = go.Figure()
    wd_range = list(range(50, 2001, 10))

    for f in focal_lengths:
        f_line = [w * sensor_w / f for w in wd_range]
        fig.add_trace(
            go.Scatter(
                x=f_line,
                y=wd_range,
                name=f"{f}mm",
                line=dict(width=4 if f == focal_len else 1.5)
            )
        )

    fig.add_shape(
        type="line",
        x0=0, y0=wd_val, x1=fov_w, y1=wd_val,
        line=dict(color="Red", dash="dash")
    )
    fig.add_shape(
        type="line",
        x0=fov_w, y0=0, x1=fov_w, y1=wd_val,
        line=dict(color="Red", dash="dash")
    )

    fig.update_layout(
        xaxis_title="FOV Width (mm)",
        yaxis_title="WD (mm)",
        height=500,
        margin=dict(t=20, b=20)
    )

    st.plotly_chart(fig, use_container_width=True)

    m1, m2, m3 = st.columns(3)
    m1.metric("Current WD", f"{wd_val:.1f} mm")
    m2.metric("FOV Width", f"{fov_w:.1f} mm")
    m3.metric("Resolution", f"{(fov_w / cam['w'] * 1000):.2f} um/px")

    st.write("---")

    dof = (2 * f_value * (cam["px"] * 2 / 1000) * (wd_val ** 2)) / (focal_len ** 2)
    st.info(f"DOF: ±{dof / 2:.2f} mm (Total {dof:.2f} mm)")